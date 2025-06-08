#!/usr/bin/env python3
"""
Multi-format document processor with OCR for intelligence analysis
Supports: TXT, PDF, DOCX, PPTX, Images with OCR capability
"""

import os
import io
from typing import Dict, Tuple, Optional
import logging

# Document processing libraries
try:
    import PyPDF2
    import pdfplumber

    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document as DocxDocument

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl

    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

# OCR libraries
try:
    import pytesseract
    from PIL import Image
    import pdf2image

    OCR_AVAILABLE = True
    print("✅ OCR libraries available (Tesseract + PIL + pdf2image)")
except ImportError as e:
    OCR_AVAILABLE = False
    print(f"⚠️ OCR libraries not available: {e}")


class DocumentProcessor:
    """Handles extraction of text from various document formats with OCR support"""

    def __init__(self):
        self.supported_formats = self._get_supported_formats()
        print(f"📄 Document processor initialized with OCR support")
        print(f"📁 Supported formats: {list(self.supported_formats.keys())}")

        # Test OCR availability
        if OCR_AVAILABLE:
            self._test_ocr_setup()

    def _test_ocr_setup(self):
        """Test if OCR is properly configured"""
        try:
            # Test Tesseract
            version = pytesseract.get_tesseract_version()
            print(f"✅ Tesseract OCR version: {version}")
        except Exception as e:
            print(f"⚠️ Tesseract OCR test failed: {e}")
            print("💡 Make sure Tesseract is installed: brew install tesseract")

    def _get_supported_formats(self) -> Dict[str, bool]:
        """Get dictionary of supported file formats"""
        formats = {
            '.txt': True,
            '.pdf': PDF_AVAILABLE,
            '.docx': DOCX_AVAILABLE,
            '.doc': DOCX_AVAILABLE,
            '.pptx': PPTX_AVAILABLE,
            '.xlsx': EXCEL_AVAILABLE,
            '.xls': EXCEL_AVAILABLE
        }

        # Add image formats if OCR is available
        if OCR_AVAILABLE:
            image_formats = {
                '.jpg': True,
                '.jpeg': True,
                '.png': True,
                '.tiff': True,
                '.tif': True,
                '.bmp': True,
                '.gif': True
            }
            formats.update(image_formats)

        return formats

    def get_file_info(self, file_obj) -> Tuple[str, str, int]:
        """Extract file information from uploaded file object"""
        filename = file_obj.name
        file_extension = os.path.splitext(filename)[1].lower()
        file_size = len(file_obj.getvalue()) if hasattr(file_obj, 'getvalue') else 0

        return filename, file_extension, file_size

    def is_supported(self, file_extension: str) -> bool:
        """Check if file format is supported"""
        return file_extension.lower() in self.supported_formats and self.supported_formats[file_extension.lower()]

    def extract_text(self, file_obj, filename: str = None) -> Tuple[str, Dict]:
        """
        Extract text from uploaded file object with OCR support
        Returns: (extracted_text, metadata)
        """
        if filename is None:
            filename = getattr(file_obj, 'name', 'unknown')

        file_extension = os.path.splitext(filename)[1].lower()

        metadata = {
            'filename': filename,
            'file_type': file_extension,
            'extraction_method': 'unknown',
            'page_count': 0,
            'char_count': 0,
            'word_count': 0,
            'ocr_used': False,
            'errors': []
        }

        try:
            if file_extension == '.txt':
                text, metadata = self._extract_txt(file_obj, metadata)
            elif file_extension == '.pdf':
                text, metadata = self._extract_pdf_smart(file_obj, metadata)
            elif file_extension in ['.docx', '.doc']:
                text, metadata = self._extract_docx(file_obj, metadata)
            elif file_extension == '.pptx':
                text, metadata = self._extract_pptx(file_obj, metadata)
            elif file_extension in ['.xlsx', '.xls']:
                text, metadata = self._extract_excel(file_obj, metadata)
            elif file_extension in ['.jpg', '.jpeg', '.png', '.tiff', '.tif', '.bmp', '.gif']:
                text, metadata = self._extract_image_ocr(file_obj, metadata)
            else:
                raise ValueError(f"Unsupported file format: {file_extension}")

            # Update metadata
            metadata['char_count'] = len(text)
            metadata['word_count'] = len(text.split())

            return text, metadata

        except Exception as e:
            error_msg = f"Error extracting text from {filename}: {str(e)}"
            metadata['errors'].append(error_msg)
            print(f"❌ {error_msg}")
            return "", metadata

    def _extract_txt(self, file_obj, metadata: Dict) -> Tuple[str, Dict]:
        """Extract text from TXT files"""
        metadata['extraction_method'] = 'direct_read'

        try:
            # Try UTF-8 first
            content = file_obj.read().decode('utf-8')
        except UnicodeDecodeError:
            # Fallback to latin-1
            file_obj.seek(0)
            content = file_obj.read().decode('latin-1', errors='ignore')
            metadata['errors'].append("Used latin-1 encoding fallback")

        return content, metadata

    def _extract_pdf_smart(self, file_obj, metadata: Dict) -> Tuple[str, Dict]:
        """Smart PDF extraction: try regular methods first, then OCR"""
        if not PDF_AVAILABLE:
            raise ImportError("PDF processing libraries not available")

        # First, try regular PDF text extraction
        try:
            text, metadata = self._extract_pdf_regular(file_obj, metadata)

            # Check if we got meaningful text (more than just whitespace/gibberish)
            meaningful_text = ''.join(c for c in text if c.isalnum())

            if len(meaningful_text) > 50:  # Got good text
                return text, metadata
            else:
                # Not enough meaningful text, try OCR
                metadata['errors'].append(
                    f"Regular extraction got limited text ({len(meaningful_text)} chars), trying OCR")
                raise ValueError("Limited text extracted, trying OCR")

        except Exception as e:
            # Regular extraction failed or got poor results, try OCR
            if OCR_AVAILABLE:
                metadata['errors'].append(f"Regular PDF extraction issue: {str(e)}")
                file_obj.seek(0)
                return self._extract_pdf_ocr(file_obj, metadata)
            else:
                raise ValueError(f"PDF extraction failed and OCR not available: {str(e)}")

    def _extract_pdf_regular(self, file_obj, metadata: Dict) -> Tuple[str, Dict]:
        """Extract text from text-based PDFs using pdfplumber/PyPDF2"""
        text = ""

        # Method 1: Try pdfplumber (better for complex layouts)
        try:
            import pdfplumber
            with pdfplumber.open(file_obj) as pdf:
                metadata['page_count'] = len(pdf.pages)
                metadata['extraction_method'] = 'pdfplumber'

                for page_num, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text += f"\n--- Page {page_num + 1} ---\n"
                        text += page_text

                if text.strip():
                    return text, metadata
        except Exception as e:
            metadata['errors'].append(f"pdfplumber failed: {str(e)}")

        # Method 2: Fallback to PyPDF2
        try:
            file_obj.seek(0)
            pdf_reader = PyPDF2.PdfReader(file_obj)
            metadata['page_count'] = len(pdf_reader.pages)
            metadata['extraction_method'] = 'PyPDF2'

            for page_num, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text += f"\n--- Page {page_num + 1} ---\n"
                    text += page_text

            return text, metadata

        except Exception as e:
            metadata['errors'].append(f"PyPDF2 failed: {str(e)}")
            raise ValueError("All regular PDF extraction methods failed")

    def _extract_pdf_ocr(self, file_obj, metadata: Dict) -> Tuple[str, Dict]:
        """Extract text from PDF using OCR (for scanned documents)"""
        if not OCR_AVAILABLE:
            raise ImportError("OCR libraries not available")

        metadata['extraction_method'] = 'OCR (PDF to images)'
        metadata['ocr_used'] = True
        text = ""

        try:
            # Convert PDF pages to images
            print("🔄 Converting PDF to images for OCR...")
            images = pdf2image.convert_from_bytes(file_obj.read(), dpi=200)
            metadata['page_count'] = len(images)

            for page_num, image in enumerate(images):
                print(f"🔍 OCR processing page {page_num + 1}/{len(images)}...")

                # Perform OCR on each page
                page_text = pytesseract.image_to_string(
                    image,
                    lang='eng',
                    config='--psm 1 --oem 3'  # Page segmentation mode and OCR engine mode
                )

                if page_text.strip():
                    text += f"\n--- Page {page_num + 1} (OCR) ---\n"
                    text += page_text

            print(f"✅ OCR completed: {len(text)} characters extracted")
            return text, metadata

        except Exception as e:
            metadata['errors'].append(f"OCR extraction failed: {str(e)}")
            raise ValueError(f"OCR extraction failed: {str(e)}")

    def _extract_image_ocr(self, file_obj, metadata: Dict) -> Tuple[str, Dict]:
        """Extract text from image files using OCR"""
        if not OCR_AVAILABLE:
            raise ImportError("OCR libraries not available")

        metadata['extraction_method'] = 'OCR (Image)'
        metadata['ocr_used'] = True
        metadata['page_count'] = 1

        try:
            print("🔍 Performing OCR on image...")

            # Open image and perform OCR
            image = Image.open(file_obj)

            # Convert to RGB if necessary (for better OCR)
            if image.mode in ('RGBA', 'LA', 'P'):
                image = image.convert('RGB')

            # Perform OCR with optimized settings
            text = pytesseract.image_to_string(
                image,
                lang='eng',
                config='--psm 1 --oem 3'
            )

            print(f"✅ Image OCR completed: {len(text)} characters extracted")
            return text, metadata

        except Exception as e:
            raise ValueError(f"Failed to extract from image: {str(e)}")

    def _extract_docx(self, file_obj, metadata: Dict) -> Tuple[str, Dict]:
        """Extract text from DOCX files"""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx library not available")

        metadata['extraction_method'] = 'python-docx'

        try:
            doc = DocxDocument(file_obj)
            text = ""

            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"

            # Extract tables
            for table in doc.tables:
                text += "\n--- Table ---\n"
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    text += " | ".join(row_text) + "\n"

            metadata['page_count'] = 1
            return text, metadata

        except Exception as e:
            raise ValueError(f"Failed to extract from DOCX: {str(e)}")

    def _extract_pptx(self, file_obj, metadata: Dict) -> Tuple[str, Dict]:
        """Extract text from PowerPoint files"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx library not available")

        metadata['extraction_method'] = 'python-pptx'

        try:
            prs = Presentation(file_obj)
            text = ""

            for slide_num, slide in enumerate(prs.slides):
                text += f"\n--- Slide {slide_num + 1} ---\n"

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"

            metadata['page_count'] = len(prs.slides)
            return text, metadata

        except Exception as e:
            raise ValueError(f"Failed to extract from PPTX: {str(e)}")

    def _extract_excel(self, file_obj, metadata: Dict) -> Tuple[str, Dict]:
        """Extract text from Excel files"""
        if not EXCEL_AVAILABLE:
            raise ImportError("openpyxl library not available")

        metadata['extraction_method'] = 'openpyxl'

        try:
            workbook = openpyxl.load_workbook(file_obj)
            text = ""

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"\n--- Sheet: {sheet_name} ---\n"

                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            row_text.append(str(cell))
                    if any(row_text):
                        text += " | ".join(row_text) + "\n"

            metadata['page_count'] = len(workbook.sheetnames)
            return text, metadata

        except Exception as e:
            raise ValueError(f"Failed to extract from Excel: {str(e)}")

    def get_missing_dependencies(self) -> Dict[str, str]:
        """Get list of missing dependencies for unsupported formats"""
        missing = {}

        if not PDF_AVAILABLE:
            missing['PDF'] = "pip install PyPDF2 pdfplumber"
        if not DOCX_AVAILABLE:
            missing['DOCX'] = "pip install python-docx"
        if not PPTX_AVAILABLE:
            missing['PPTX'] = "pip install python-pptx"
        if not EXCEL_AVAILABLE:
            missing['Excel'] = "pip install openpyxl"
        if not OCR_AVAILABLE:
            missing['OCR/Images'] = "brew install tesseract && pip install pytesseract Pillow pdf2image"

        return missing


# Testing function
if __name__ == "__main__":
    processor = DocumentProcessor()

    print("🧪 Document Processor with OCR Test")
    print(f"📁 Supported formats: {list(processor.supported_formats.keys())}")

    missing = processor.get_missing_dependencies()
    if missing:
        print("\n⚠️ Missing dependencies:")
        for format_name, install_cmd in missing.items():
            print(f"  {format_name}: {install_cmd}")
    else:
        print("\n✅ All document processing libraries available!")

    if OCR_AVAILABLE:
        print("🔍 OCR capabilities: Ready for scanned documents and images")
    else:
        print("❌ OCR not available - install Tesseract and OCR libraries")